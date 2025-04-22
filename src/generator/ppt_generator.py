from pptx import Presentation

def generate_ppt(template_path: str, output_path: str, fields: dict) -> str:
    """
    - template_path: ruta al .pptx con placeholders {{campo}}
    - output_path: ruta donde se guardará el PPT final
    - fields: diccionario con clave=campo, valor=texto a pegar
    """
    prs = Presentation(template_path)

    for slide in prs.slides:
        # Reemplazo en cuadros de texto
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        for key, val in fields.items():
                            placeholder = f"{{{{{key}}}}}"   # resultado: {{nombre_empresa}}
                            if placeholder in text:
                                text = text.replace(placeholder, str(val))
                        run.text = text

        # Reemplazo en tablas 
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text
                    for key, val in fields.items():
                        ph = f"{{{{{key}}}}}"
                        if ph in cell_text:
                            cell.text = cell_text.replace(ph, str(val))

    prs.save(output_path)
    return output_path


